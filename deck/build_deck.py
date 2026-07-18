"""
Builds agentic-erp-study.pptx, an 11 slide deck presenting the findings of
STUDY.md and findings/results.md (agentic-erp-demo repo).

Design: clean consultant style, 16:9, white ground, one dark ink color
(near black navy), one accent (deep teal blue, not SAP branding), large
numbers for stats, simple native shapes for diagrams. No clip art.

Run:  .venv/bin/python deck/build_deck.py
Out:  deck/agentic-erp-study.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
import copy

# ---------------------------------------------------------------------------
# Palette and type
# ---------------------------------------------------------------------------

INK = RGBColor(0x0B, 0x14, 0x2A)        # near black navy
ACCENT = RGBColor(0x0E, 0x6B, 0x6B)     # deep teal
ACCENT_LIGHT = RGBColor(0xE4, 0xF0, 0xEF)  # pale teal tint for fills
PAPER = RGBColor(0xFF, 0xFF, 0xFF)      # white ground
GREY = RGBColor(0x6B, 0x72, 0x80)       # muted supporting text
GREY_LINE = RGBColor(0xDD, 0xE1, 0xE6)  # hairlines
RED_MUTED = RGBColor(0x9C, 0x3B, 0x2E)  # for the blocked / failure marker

FONT_HEAD = "Georgia"
FONT_BODY = "Helvetica Neue"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.7)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def add_slide():
    return prs.slides.add_slide(BLANK)


def set_bg(slide, color=PAPER):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def textbox(slide, left, top, width, height, text, size=18, color=INK,
            bold=False, italic=False, font=FONT_BODY, align=PP_ALIGN.LEFT,
            anchor=MSO_ANCHOR.TOP, line_spacing=1.0, space_after=0):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.italic = italic
        r.font.name = font
    return box


def add_multirun(slide, left, top, width, height, runs, size=18,
                  align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0):
    """runs: list of (text, color, bold, italic, font, size_override)"""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for text, color, bold, italic, font, sz in runs:
        r = p.add_run()
        r.text = text
        r.font.size = Pt(sz if sz else size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.italic = italic
        r.font.name = font
    return box


def rect(slide, left, top, width, height, fill_color=None, line_color=None,
         line_width=None, shadow=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.shadow.inherit = False
    if fill_color is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill_color
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = line_width or Pt(0.75)
    return shp


def rounded_rect(slide, left, top, width, height, fill_color=None,
                  line_color=None, line_width=None, radius=0.06):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.shadow.inherit = False
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    if fill_color is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill_color
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = line_width or Pt(0.75)
    return shp


def hline(slide, left, top, width, color=GREY_LINE, weight=Pt(1)):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left, top, left + width, top)
    ln.line.color.rgb = color
    ln.line.width = weight
    ln.shadow.inherit = False
    return ln


def vline(slide, left, top, height, color=GREY_LINE, weight=Pt(1)):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left, top, left, top + height)
    ln.line.color.rgb = color
    ln.line.width = weight
    ln.shadow.inherit = False
    return ln


def arrow(slide, left, top, width, height, color=ACCENT, weight=Pt(1.5)):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left, top, left + width, top + height)
    ln.line.color.rgb = color
    ln.line.width = weight
    ln.shadow.inherit = False
    line_elem = ln.line._get_or_add_ln()
    tail = line_elem.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    line_elem.append(tail)
    return ln


SLIDE_NUM = [0]


def kicker_and_title(slide, kicker, title, title_size=30):
    """Standard header block used on every content slide (2-10)."""
    textbox(slide, MARGIN, Inches(0.5), Inches(8), Inches(0.35),
            kicker.upper(), size=12.5, color=ACCENT, bold=True,
            font=FONT_BODY)
    textbox(slide, MARGIN, Inches(0.82), Inches(11.9), Inches(0.9),
            title, size=title_size, color=INK, bold=True, font=FONT_HEAD)
    hline(slide, MARGIN, Inches(1.62), SLIDE_W - 2 * MARGIN, color=GREY_LINE, weight=Pt(1))


def footer(slide, n):
    textbox(slide, MARGIN, Inches(7.1), Inches(6), Inches(0.3),
            "Can an AI agent operate an ERP? — July 2026 study", size=9.5,
            color=GREY, font=FONT_BODY)
    textbox(slide, SLIDE_W - MARGIN - Inches(1.5), Inches(7.1), Inches(1.5), Inches(0.3),
            f"{n:02d} / 11", size=9.5, color=GREY, font=FONT_BODY, align=PP_ALIGN.RIGHT)


def stat_block(slide, left, top, width, number, label, number_size=44, label_size=12.5,
               number_color=ACCENT):
    textbox(slide, left, top, width, Inches(0.75), number, size=number_size,
            color=number_color, bold=True, font=FONT_HEAD, align=PP_ALIGN.LEFT)
    label_top = top + Inches(0.72 * max(1.0, number_size / 44))
    textbox(slide, left, label_top, width, Inches(0.6), label, size=label_size,
            color=GREY, font=FONT_BODY, align=PP_ALIGN.LEFT, line_spacing=1.15)


# ===========================================================================
# Slide 1 — Title
# ===========================================================================
s = add_slide()
set_bg(s)
SLIDE_NUM[0] += 1

# accent rule at top
rect(s, 0, 0, SLIDE_W, Inches(0.12), fill_color=ACCENT)

textbox(s, MARGIN, Inches(2.55), Inches(11.9), Inches(0.4),
        "A HANDS ON STUDY • JULY 2026", size=13, color=ACCENT, bold=True,
        font=FONT_BODY)

textbox(s, MARGIN, Inches(2.95), Inches(11.9), Inches(1.9),
        "Can an AI agent operate an ERP?", size=46, color=INK, bold=True,
        font=FONT_HEAD, line_spacing=1.05)

textbox(s, MARGIN, Inches(4.35), Inches(11.0), Inches(0.7),
        "A hands on study against SAP's public APIs, July 2026",
        size=18, color=GREY, italic=True, font=FONT_BODY)

hline(s, MARGIN, Inches(5.3), Inches(4.2), color=GREY_LINE, weight=Pt(1))

textbox(s, MARGIN, Inches(6.75), Inches(9), Inches(0.4),
        "github.com/asish-singh/agentic-erp-demo", size=13.5, color=INK,
        bold=True, font=FONT_BODY)

# ===========================================================================
# Slide 2 — Why now
# ===========================================================================
s = add_slide()
set_bg(s)
kicker_and_title(s, "Why now", "SAP is telling the market agents will run the enterprise")

body_left = MARGIN
col_w = Inches(6.9)

textbox(s, body_left, Inches(2.0), col_w, Inches(2.6),
        "In May 2026 SAP announced its Autonomous Enterprise push, "
        "more than 200 specialized agents planned across its product line. "
        "The pitch is simple. AI agents will operate enterprise software, "
        "not just assist the people who use it.",
        size=16, color=INK, font=FONT_BODY, line_spacing=1.3)

textbox(s, body_left, Inches(4.1), col_w, Inches(1.9),
        "Nobody outside SAP has published what happens when an ordinary "
        "agent, with no special access, tries this against the public APIs "
        "every customer and partner can already reach. This study asks that "
        "question and shows its work.",
        size=16, color=INK, font=FONT_BODY, line_spacing=1.3)

# right stat card
card_left = Inches(8.35)
card_w = Inches(4.25)
rounded_rect(s, card_left, Inches(2.05), card_w, Inches(3.9), fill_color=ACCENT_LIGHT, radius=0.05)
stat_block(s, card_left + Inches(0.4), Inches(2.45), card_w - Inches(0.8),
           "200+", "agents announced by SAP\nacross its product line, May 2026",
           number_size=52)
hline(s, card_left + Inches(0.4), Inches(4.15), card_w - Inches(0.8), color=RGBColor(0xC7,0xDA,0xD9))
textbox(s, card_left + Inches(0.4), Inches(4.35), card_w - Inches(0.8), Inches(1.3),
        "The question this study tests\nis narrower and it is not\ncurrently answered in public",
        size=14, color=INK, font=FONT_BODY, line_spacing=1.25)

footer(s, 2)

# ===========================================================================
# Slide 3 — Method
# ===========================================================================
s = add_slide()
set_bg(s)
kicker_and_title(s, "Method", "One agent loop, a free model, SAP's public sandbox")

diagram_top = Inches(2.15)
box_h = Inches(1.05)
box_w = Inches(2.55)
gap = Inches(0.55)
y = diagram_top

box1_l = MARGIN
box2_l = box1_l + box_w + gap
box3_l = box2_l + box_w + gap

# Box 1: instruction
rounded_rect(s, box1_l, y, box_w, box_h, fill_color=PAPER, line_color=INK, line_width=Pt(1.25), radius=0.08)
textbox(s, box1_l, y + Inches(0.12), box_w, Inches(0.35), "Plain language",
        size=11.5, color=GREY, align=PP_ALIGN.CENTER, font=FONT_BODY)
textbox(s, box1_l, y + Inches(0.42), box_w, Inches(0.5), "Instruction",
        size=16.5, color=INK, bold=True, align=PP_ALIGN.CENTER, font=FONT_HEAD)

# Box 2: agent loop with 4 tools
rounded_rect(s, box2_l, y, box_w, box_h, fill_color=ACCENT, radius=0.08)
textbox(s, box2_l, y + Inches(0.1), box_w, Inches(0.3), "Agent loop • 4 tools • 15 turns max",
        size=10.5, color=RGBColor(0xE6,0xF3,0xF2), align=PP_ALIGN.CENTER, font=FONT_BODY)
textbox(s, box2_l, y + Inches(0.4), box_w, Inches(0.5), "GPT-4o mini",
        size=16.5, color=PAPER, bold=True, align=PP_ALIGN.CENTER, font=FONT_HEAD)

# Box 3: SAP sandbox
rounded_rect(s, box3_l, y, box_w, box_h, fill_color=PAPER, line_color=INK, line_width=Pt(1.25), radius=0.08)
textbox(s, box3_l, y + Inches(0.12), box_w, Inches(0.35), "S/4HANA Cloud",
        size=11.5, color=GREY, align=PP_ALIGN.CENTER, font=FONT_BODY)
textbox(s, box3_l, y + Inches(0.42), box_w, Inches(0.5), "SAP sandbox APIs",
        size=16.5, color=INK, bold=True, align=PP_ALIGN.CENTER, font=FONT_HEAD)

# arrows
arrow(s, box1_l + box_w, y + box_h/2, gap, 0)
arrow(s, box2_l + box_w, y + box_h/2, gap, 0)

# tools list under the agent box
tools_top = y + box_h + Inches(0.28)
textbox(s, box2_l, tools_top, box_w, Inches(1.3),
        "list services\nread from an API\nwrite to an API\nfinish",
        size=12.5, color=INK, font=FONT_BODY, align=PP_ALIGN.CENTER, line_spacing=1.35)

# "all logged" bracket under whole diagram
log_top = y + box_h + Inches(1.75)
hline(s, box1_l, log_top, box3_l + box_w - box1_l, color=INK, weight=Pt(1))
textbox(s, box1_l, log_top + Inches(0.12), box3_l + box_w - box1_l, Inches(0.4),
        "Every model response and every API call logged in full to JSONL",
        size=13, color=INK, italic=True, align=PP_ALIGN.CENTER, font=FONT_BODY)

# zero cost stack strip
strip_top = Inches(6.0)
rounded_rect(s, MARGIN, strip_top, SLIDE_W - 2*MARGIN, Inches(0.85), fill_color=ACCENT_LIGHT, radius=0.12)
textbox(s, MARGIN + Inches(0.35), strip_top + Inches(0.14), Inches(2.2), Inches(0.55),
        "Zero cost stack", size=14, color=ACCENT, bold=True, font=FONT_BODY,
        anchor=MSO_ANCHOR.MIDDLE)
textbox(s, MARGIN + Inches(2.7), strip_top + Inches(0.14), Inches(9.8), Inches(0.55),
        "Free model tier, free SAP sandbox key, runner executes locally or in GitHub Actions",
        size=14, color=INK, font=FONT_BODY, anchor=MSO_ANCHOR.MIDDLE)

footer(s, 3)

# ===========================================================================
# Slide 4 — The five tasks
# ===========================================================================
s = add_slide()
set_bg(s)
kicker_and_title(s, "The tasks", "Five ordinary shapes of ERP work")

tasks = [
    ("01", "Find a master data record", "Find a supplier by name and report its account group"),
    ("02", "Filter transactional documents", "List purchase orders for a supplier and report currencies"),
    ("03", "Create a document", "Create a purchase order"),
    ("04", "Summarize documents", "Summarize a supplier's invoices"),
    ("05", "Cross reference two objects", "Report a product's base unit and type"),
]

row_top = Inches(2.15)
row_h = Inches(0.92)
for i, (num, kind, desc) in enumerate(tasks):
    top = row_top + i * row_h
    if i > 0:
        hline(s, MARGIN, top, SLIDE_W - 2*MARGIN, color=GREY_LINE, weight=Pt(0.75))
    textbox(s, MARGIN, top + Inches(0.14), Inches(0.9), Inches(0.6), num,
            size=22, color=ACCENT, bold=True, font=FONT_HEAD)
    textbox(s, MARGIN + Inches(1.05), top + Inches(0.1), Inches(3.9), Inches(0.7), kind,
            size=15.5, color=INK, bold=True, font=FONT_BODY, line_spacing=1.1)
    textbox(s, MARGIN + Inches(5.15), top + Inches(0.1), Inches(6.9), Inches(0.7), desc,
            size=15.5, color=GREY, font=FONT_BODY, line_spacing=1.1)

footer(s, 4)

# ===========================================================================
# Slide 5 — Results table
# ===========================================================================
s = add_slide()
set_bg(s)
kicker_and_title(s, "Results", "Four of five tasks succeeded")

# big stat
stat_block(s, MARGIN, Inches(1.95), Inches(3.4), "4 / 5", "tasks succeeded\non first evaluation",
            number_size=46)

table_left = Inches(4.7)
table_top = Inches(2.0)
table_w = SLIDE_W - MARGIN - table_left

results = [
    ("Find a supplier and report its account group", "Success", "13", "11"),
    ("List purchase orders for a supplier and report currencies", "Success", "3", "1"),
    ("Create a purchase order", "Blocked by platform", "4", "2"),
    ("Summarize supplier invoices", "Success", "6", "4"),
    ("Report a product's base unit and type", "Success", "6", "4"),
]

col_task_w = Inches(4.9)
col_out_w = Inches(1.9)
col_turn_w = Inches(0.95)
col_api_w = Inches(1.0)

headers = ["Task", "Outcome", "Turns", "API calls"]
col_lefts = [table_left, table_left+col_task_w, table_left+col_task_w+col_out_w,
             table_left+col_task_w+col_out_w+col_turn_w]
col_widths = [col_task_w, col_out_w, col_turn_w, col_api_w]

hline(s, table_left, table_top, table_w, color=INK, weight=Pt(1.25))
for i, h in enumerate(headers):
    align = PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER
    textbox(s, col_lefts[i], table_top + Inches(0.08), col_widths[i], Inches(0.35), h,
            size=12, color=GREY, bold=True, font=FONT_BODY, align=align)
hline(s, table_left, table_top + Inches(0.5), table_w, color=INK, weight=Pt(1))

row_h5 = Inches(0.72)
for r, (task, outcome, turns, api) in enumerate(results):
    top = table_top + Inches(0.5) + r * row_h5
    textbox(s, col_lefts[0], top + Inches(0.16), col_widths[0], Inches(0.6), task,
            size=13.5, color=INK, font=FONT_BODY, line_spacing=1.1)
    is_blocked = outcome != "Success"
    oc = RED_MUTED if is_blocked else ACCENT
    textbox(s, col_lefts[1], top + Inches(0.2), col_widths[1], Inches(0.4), outcome,
            size=13, color=oc, bold=True, font=FONT_BODY, align=PP_ALIGN.CENTER)
    textbox(s, col_lefts[2], top + Inches(0.2), col_widths[2], Inches(0.4), turns,
            size=13.5, color=INK, font=FONT_BODY, align=PP_ALIGN.CENTER)
    textbox(s, col_lefts[3], top + Inches(0.2), col_widths[3], Inches(0.4), api,
            size=13.5, color=INK, font=FONT_BODY, align=PP_ALIGN.CENTER)
    hline(s, table_left, top + row_h5, table_w, color=GREY_LINE, weight=Pt(0.6))

textbox(s, MARGIN, Inches(6.35), Inches(6.8), Inches(0.6),
        "The fifth task was blocked by platform policy, not agent failure",
        size=13, color=GREY, italic=True, font=FONT_BODY)

footer(s, 5)

# ===========================================================================
# Slide 6 — Does a smarter brain fix it?
# ===========================================================================
s = add_slide()
set_bg(s)
kicker_and_title(s, "Model comparison", "Does a smarter brain fix it?")

textbox(s, MARGIN, Inches(1.85), Inches(11.9), Inches(0.6),
        "The same five tasks, repeated with a frontier model and a strong "
        "open source model, same harness, same tools, same logging.",
        size=15, color=INK, font=FONT_BODY, line_spacing=1.25)

table_left6 = MARGIN
table_top6 = Inches(2.45)
table_w6 = SLIDE_W - 2 * MARGIN

col0_w = Inches(3.6)
col_model_w = (table_w6 - col0_w) / 3

headers6 = ["Task", "GPT-4o mini (small)", "GPT-4.1 (frontier)", "Llama 4 Maverick (open)"]
col_lefts6 = [table_left6, table_left6 + col0_w, table_left6 + col0_w + col_model_w,
              table_left6 + col0_w + 2 * col_model_w]
col_widths6 = [col0_w, col_model_w, col_model_w, col_model_w]

hline(s, table_left6, table_top6, table_w6, color=INK, weight=Pt(1.25))
for i, h in enumerate(headers6):
    align = PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER
    textbox(s, col_lefts6[i], table_top6 + Inches(0.08), col_widths6[i], Inches(0.45), h,
            size=12, color=GREY, bold=True, font=FONT_BODY, align=align, line_spacing=1.05)
hline(s, table_left6, table_top6 + Inches(0.6), table_w6, color=INK, weight=Pt(1))

rows6 = [
    ("Find a supplier", "success, 11 calls", "success, 1 call", "gave up, 3 calls"),
    ("List purchase orders", "success", "success", "success"),
    ("Create a purchase order", "blocked, 405", "blocked, 405", "blocked, 405"),
    ("Summarize invoices", "success", "success", "success"),
    ("Product base unit and type", "success", "success", "success"),
]

row_h6 = Inches(0.48)
for r, (task, a, b, c) in enumerate(rows6):
    top = table_top6 + Inches(0.6) + r * row_h6
    textbox(s, col_lefts6[0], top + Inches(0.12), col_widths6[0], Inches(0.4), task,
            size=13, color=INK, font=FONT_BODY)
    is_blocked = task == "Create a purchase order"
    is_giveup = "gave up" in b or "gave up" in c
    for idx, val in enumerate((a, b, c)):
        cell_color = RED_MUTED if is_blocked else (RED_MUTED if "gave up" in val else INK)
        textbox(s, col_lefts6[idx + 1], top + Inches(0.12), col_widths6[idx + 1], Inches(0.4),
                val, size=13, color=cell_color, font=FONT_BODY, align=PP_ALIGN.CENTER)
    hline(s, table_left6, top + row_h6, table_w6, color=GREY_LINE, weight=Pt(0.6))

strip_top6 = table_top6 + Inches(0.6) + len(rows6) * row_h6 + Inches(0.2)
strip_h6 = Inches(0.8)
rounded_rect(s, table_left6, strip_top6, table_w6, strip_h6, fill_color=INK, radius=0.08)
textbox(s, table_left6 + Inches(0.4), strip_top6, table_w6 - Inches(0.8), strip_h6,
        "Model quality buys efficiency, not new capability. No model gets through a "
        "platform policy. All three hit the same 405 wall on the write task.",
        size=14.5, color=PAPER, italic=True, font=FONT_BODY, line_spacing=1.25,
        anchor=MSO_ANCHOR.MIDDLE)

footer(s, 6)

# ===========================================================================
# Slide 7 — Finding 1
# ===========================================================================
s = add_slide()
set_bg(s)
kicker_and_title(s, "Finding 1", "Reading works, but discovery is expensive")

stat_block(s, MARGIN, Inches(2.0), Inches(3.7), "11", "API calls to find\none supplier by name",
            number_size=52)

textbox(s, MARGIN, Inches(4.0), Inches(3.7), Inches(2.6),
        "The agent completed every read task. But the simplest sounding one, "
        "find one supplier by name, took 11 calls because the agent had to "
        "discover SAP's data model by trial and error.",
        size=14.5, color=INK, font=FONT_BODY, line_spacing=1.3)

# right column: entity guessing example
right_l = Inches(5.4)
right_w = SLIDE_W - MARGIN - right_l
rounded_rect(s, right_l, Inches(2.0), right_w, Inches(3.9), fill_color=ACCENT_LIGHT, radius=0.05)
textbox(s, right_l+Inches(0.4), Inches(2.3), right_w-Inches(0.8), Inches(0.4),
        "Entity names the agent guessed and SAP rejected", size=13.5, color=ACCENT,
        bold=True, font=FONT_BODY)
textbox(s, right_l+Inches(0.4), Inches(2.85), right_w-Inches(0.8), Inches(0.9),
        "A_PaymentTerms\nA_SupplierRole", size=17, color=INK, bold=True,
        font="Courier New", line_spacing=1.4)
hline(s, right_l+Inches(0.4), Inches(3.9), right_w-Inches(0.8), color=RGBColor(0xC7,0xDA,0xD9))
textbox(s, right_l+Inches(0.4), Inches(4.1), right_w-Inches(0.8), Inches(1.6),
        "It also used a filter function the OData v2 gateway rejects outright, "
        "and recovered each time by reading the error and adjusting course. "
        "The stumbles are not model stupidity, they are the cost of a fifty "
        "year old data model exposed through an API that assumes the caller "
        "already knows it.",
        size=13, color=INK, font=FONT_BODY, line_spacing=1.28)

footer(s, 7)

# ===========================================================================
# Slide 8 — Finding 2
# ===========================================================================
s = add_slide()
set_bg(s)
kicker_and_title(s, "Finding 2", "Writing to the enterprise is walled off")

stat_block(s, MARGIN, Inches(2.0), Inches(3.7), "405", "HTTP response to the\npurchase order write attempt",
            number_size=52, number_color=RED_MUTED)

textbox(s, MARGIN, Inches(4.05), Inches(3.7), Inches(2.5),
        "SAP's message was direct, the public sandbox supports GET operations "
        "only, and write operations must be tested against a customer's own "
        "system.",
        size=14.5, color=INK, font=FONT_BODY, line_spacing=1.3)

right_l = Inches(5.4)
right_w = SLIDE_W - MARGIN - right_l
rounded_rect(s, right_l, Inches(2.0), right_w, Inches(1.55), fill_color=INK, radius=0.06)
textbox(s, right_l+Inches(0.4), Inches(2.18), right_w-Inches(0.8), Inches(1.2),
        "“The public sandbox supports GET operations only. Write "
        "operations must be tested against a customer's own system.”",
        size=15, color=PAPER, italic=True, font=FONT_BODY, line_spacing=1.35,
        anchor=MSO_ANCHOR.MIDDLE)

textbox(s, right_l, Inches(3.85), right_w, Inches(1.9),
        "The public evaluation surface of the world's largest ERP vendor "
        "lets an agent look but not act. Anyone claiming their agent "
        "“works with SAP” on the basis of public APIs is describing a "
        "read only integration, a distinction buyers should press on.",
        size=14, color=INK, font=FONT_BODY, line_spacing=1.32)

footer(s, 8)

# ===========================================================================
# Slide 9 — Finding 3
# ===========================================================================
s = add_slide()
set_bg(s)
kicker_and_title(s, "Finding 3", "Verbosity is a tax on agents")

stat_block(s, MARGIN, Inches(2.0), Inches(3.7), "8,000", "token limit overflowed by\na single unfiltered API reply",
            number_size=48)

textbox(s, MARGIN, Inches(4.05), Inches(3.7), Inches(2.6),
        "SAP's raw responses are large enough that one unfiltered reply "
        "overflowed the free model's request limit and killed the first "
        "run outright. The failed log is preserved in the repository.",
        size=14.5, color=INK, font=FONT_BODY, line_spacing=1.3)

# simple bar comparison: raw response vs trimmed response
right_l = Inches(5.6)
right_w = SLIDE_W - MARGIN - right_l
bar_top = Inches(2.3)
bar_h = Inches(0.55)
max_bar_w = right_w - Inches(0.2)

# raw bar (full width, overflow)
rect(s, right_l, bar_top, max_bar_w, bar_h, fill_color=RED_MUTED)
textbox(s, right_l, bar_top - Inches(0.35), max_bar_w, Inches(0.3),
        "Raw SAP response", size=12.5, color=GREY, font=FONT_BODY)
textbox(s, right_l + Inches(0.15), bar_top, max_bar_w - Inches(0.3), bar_h,
        "overflows 8,000 token limit", size=12.5, color=PAPER, bold=True,
        font=FONT_BODY, anchor=MSO_ANCHOR.MIDDLE)

bar2_top = bar_top + Inches(1.15)
trimmed_w = Inches(2.4)
rect(s, right_l, bar2_top, trimmed_w, bar_h, fill_color=ACCENT)
textbox(s, right_l, bar2_top - Inches(0.35), max_bar_w, Inches(0.3),
        "Trimmed response used by the agent", size=12.5, color=GREY, font=FONT_BODY)
textbox(s, right_l + Inches(0.15), bar2_top, trimmed_w - Inches(0.3), bar_h,
        "fits comfortably", size=12.5, color=PAPER, bold=True, font=FONT_BODY,
        anchor=MSO_ANCHOR.MIDDLE)

textbox(s, right_l, bar2_top + Inches(1.1), right_w, Inches(1.6),
        "Enterprise APIs were designed for programs that extract one field "
        "and move on. Agents read everything, so response weight becomes a "
        "real cost and a real failure mode.",
        size=13.5, color=INK, font=FONT_BODY, line_spacing=1.3)

footer(s, 9)

# ===========================================================================
# Slide 10 — Finding 4
# ===========================================================================
s = add_slide()
set_bg(s)
kicker_and_title(s, "Finding 4", "Error messages are the agent's documentation")

textbox(s, MARGIN, Inches(2.0), Inches(11.9), Inches(0.7),
        "The agent never saw SAP documentation, only error responses.",
        size=17, color=INK, font=FONT_BODY, line_spacing=1.3)

col_w9 = Inches(5.6)
gap9 = Inches(0.7)
left1 = MARGIN
left2 = left1 + col_w9 + gap9
card_top = Inches(2.9)
card_h = Inches(2.9)

rounded_rect(s, left1, card_top, col_w9, card_h, fill_color=ACCENT_LIGHT, radius=0.05)
textbox(s, left1+Inches(0.35), card_top+Inches(0.28), col_w9-Inches(0.7), Inches(0.35),
        "Specific error", size=13, color=ACCENT, bold=True, font=FONT_BODY)
textbox(s, left1+Inches(0.35), card_top+Inches(0.68), col_w9-Inches(0.7), Inches(0.85),
        "“Property contains not found in type A_BusinessPartnerType”",
        size=13.5, color=INK, italic=True, font=FONT_BODY, line_spacing=1.3)
textbox(s, left1+Inches(0.35), card_top+Inches(1.75), col_w9-Inches(0.7), Inches(0.9),
        "Result: the agent corrected course in one turn",
        size=14, color=INK, bold=True, font=FONT_BODY, line_spacing=1.3)

rounded_rect(s, left2, card_top, col_w9, card_h, fill_color=PAPER, line_color=GREY_LINE,
             line_width=Pt(1), radius=0.05)
textbox(s, left2+Inches(0.35), card_top+Inches(0.28), col_w9-Inches(0.7), Inches(0.35),
        "Vague error", size=13, color=RED_MUTED, bold=True, font=FONT_BODY)
textbox(s, left2+Inches(0.35), card_top+Inches(0.68), col_w9-Inches(0.7), Inches(0.85),
        "Generic not found or rejected responses with no field level detail",
        size=13.5, color=INK, italic=True, font=FONT_BODY, line_spacing=1.3)
textbox(s, left2+Inches(0.35), card_top+Inches(1.75), col_w9-Inches(0.7), Inches(0.9),
        "Result: repeated failing guesses",
        size=14, color=INK, bold=True, font=FONT_BODY, line_spacing=1.3)

textbox(s, MARGIN, Inches(6.15), Inches(11.9), Inches(0.7),
        "For API owners preparing for agent traffic, error message quality "
        "is no longer a developer nicety, it is the interface.",
        size=14, color=GREY, italic=True, font=FONT_BODY, line_spacing=1.3)

footer(s, 10)

# ===========================================================================
# Slide 11 — What this means
# ===========================================================================
s = add_slide()
set_bg(s)
kicker_and_title(s, "What this means", "Three implications, for three audiences")

implications = [
    ("Agent builders", "Budget for the discovery problem. Agents will need curated "
     "tool definitions or metadata access, not raw OData endpoints and optimism."),
    ("Buyers", "Ask vendors whether their SAP agent integration reads, writes, or "
     "both, and where the write path was actually tested."),
    ("Platform owners", "The levers that decide agent success are unglamorous, lean "
     "responses, precise errors, and discoverable metadata."),
]

row_top10 = Inches(2.1)
row_h10 = Inches(1.35)
for i, (aud, text) in enumerate(implications):
    top = row_top10 + i * row_h10
    textbox(s, MARGIN, top, Inches(0.7), Inches(0.6), str(i+1), size=26, color=ACCENT,
            bold=True, font=FONT_HEAD)
    textbox(s, MARGIN + Inches(0.85), top + Inches(0.02), Inches(2.6), Inches(0.9), aud,
            size=15.5, color=INK, bold=True, font=FONT_BODY, line_spacing=1.15)
    textbox(s, MARGIN + Inches(3.6), top, Inches(8.4), Inches(1.1), text,
            size=14.5, color=INK, font=FONT_BODY, line_spacing=1.3)
    if i < len(implications) - 1:
        hline(s, MARGIN, top + row_h10 - Inches(0.15), SLIDE_W - 2*MARGIN, color=GREY_LINE, weight=Pt(0.6))

rounded_rect(s, MARGIN, Inches(6.15), SLIDE_W - 2*MARGIN, Inches(0.75), fill_color=ACCENT_LIGHT, radius=0.12)
textbox(s, MARGIN + Inches(0.35), Inches(6.15) + Inches(0.2), Inches(11), Inches(0.4),
        "github.com/asish-singh/agentic-erp-demo", size=15, color=INK, bold=True,
        font=FONT_BODY, anchor=MSO_ANCHOR.MIDDLE)

footer(s, 11)

# ---------------------------------------------------------------------------
out_path = "/Users/asishsingh/Documents/Codebase/agentic-erp-demo/deck/agentic-erp-study.pptx"
prs.save(out_path)
print(f"Saved {out_path}")
